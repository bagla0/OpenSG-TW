# -*- coding: utf-8 -*-
"""windIO -> OpenSG -> RM shell/KL shell/Solid -> Timoshenko Beam flowchart as an editable PPTX."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

REPO = r"C:\Users\bagla0\OneDrive - purdue.edu\2026_195\Claude_code"
SCALE = 13.33 / 18.7                                    # matplotlib units (0..18.7) -> inches across a 13.33in slide
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(3.72)
slide = prs.slides.add_slide(prs.slide_layouts[6])     # blank

BLUE, ORANGE, GREEN = RGBColor(0xcf, 0xe3, 0xf3), RGBColor(0xfd, 0xe0, 0xc0), RGBColor(0xcf, 0xea, 0xd0)
PURPLE, YEL, TEAL = RGBColor(0xe3, 0xd4, 0xef), RGBColor(0xff, 0xf2, 0xb3), RGBColor(0xbf, 0xe6, 0xe6)
GRAY, INK = RGBColor(0x59, 0x59, 0x59), RGBColor(0x20, 0x20, 0x20)


def _xy(x, y):
    return Inches(x * SCALE), Inches((5 - y) * SCALE)


def _frame(cx, cy, w, h, fc):
    left, top = _xy(cx - w / 2, cy + h / 2)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(w * SCALE), Inches(h * SCALE))
    sp.fill.solid(); sp.fill.fore_color.rgb = fc
    sp.line.color.rgb = GRAY; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Pt(2); tf.margin_left = tf.margin_right = Pt(3)
    return sp


def _line(p, text, size, bold):
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = INK


def box(cx, cy, w, h, text, fc, size=14, bold=False):
    tf = _frame(cx, cy, w, h, fc).text_frame
    for i, ln in enumerate(text.split("\n")):
        _line(tf.paragraphs[0] if i == 0 else tf.add_paragraph(), ln, size, bold)


def proc(cx, cy, w, h, head, sub, fc):
    tf = _frame(cx, cy, w, h, fc).text_frame
    _line(tf.paragraphs[0], head, 16, True)
    for ln in sub.split("\n"):
        _line(tf.add_paragraph(), ln, 10, False)


def arrow(x0, y0, x1, y1):
    (a, b), (c, d) = _xy(x0, y0), _xy(x1, y1)
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a, b, c, d)
    cn.line.color.rgb = GRAY; cn.line.width = Pt(1.5)
    ln = cn.line._get_or_add_ln()
    ln.append(parse_xml('<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'type="triangle" w="med" len="med"/>'))


box(1.6, 2.5, 2.6, 1.1, "WindIO Blade\n(.yaml)", BLUE, 14, True)
proc(5.2, 3.8, 3.3, 1.35, "OpenSG_io", "load_blade →\nbuild_cross_section →\nemit_opensg_yaml", ORANGE)
proc(5.2, 1.2, 3.3, 1.35, "PreVABS", "emit_prevabs →\nprevabs --vabs --hm →\nconvert_sg_to_yaml", ORANGE)
box(8.9, 3.8, 2.3, 1.0, "1D shell\nSG YAML", GREEN, 13)
box(8.9, 1.2, 2.3, 1.0, "2D solid\nSG YAML", GREEN, 13)
box(11.7, 2.5, 2.0, 1.15, "OpenSG", PURPLE, 16, True)
box(14.3, 3.9, 2.0, 0.8, "RM shell", YEL, 13)
box(14.3, 2.5, 2.0, 0.8, "KL shell", YEL, 13)
box(14.3, 1.1, 2.0, 0.8, "Solid", YEL, 13)
box(16.95, 2.5, 2.7, 1.25, "Timoshenko\nBeam", TEAL, 13)

arrow(2.9, 2.8, 3.55, 3.5); arrow(2.9, 2.2, 3.55, 1.5)
arrow(6.85, 3.8, 7.75, 3.8); arrow(6.85, 1.2, 7.75, 1.2)
arrow(10.05, 3.8, 10.75, 2.85); arrow(10.05, 1.2, 10.75, 2.15)
arrow(12.7, 2.75, 13.25, 3.85); arrow(12.7, 2.5, 13.25, 2.5); arrow(12.7, 2.25, 13.25, 1.15)
arrow(15.35, 3.9, 15.55, 2.8); arrow(15.35, 2.5, 15.55, 2.5); arrow(15.35, 1.1, 15.55, 2.2)

out = os.path.join(REPO, "docs", "tutorials", "_img", "windio_workflow.pptx")
prs.save(out)
print("wrote", out)
